# from pptx import Presentation
# import os
from pptx import Presentation

# current_dir = os.path.dirname(os.path.realpath(__file__))

# print(current_dir)

prs = Presentation('template.pptx')
bullet_slidelayout = prs.slide_layouts[1]
slideContent = [['title1','point1','point1','point1'],['title2','point1','point1','point1'],['title3','point1','point1','point1']]
for eachSlide in slideContent:
    slide = prs.slides.add_slide(bullet_slidelayout)
    shapes = slide.shapes
    title_shape = shapes.title
    body_shape = shapes.placeholders[1]

    title_shape.text = eachSlide[0]
    tf = body_shape.text_frame
    for points in eachSlide[1:]:
        tf.add_paragraph().text = points
prs.save('test.pptx')